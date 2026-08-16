#!/usr/bin/env python3
"""
Builds a native PowerPoint deck from slides.md.

Every title, caption and paragraph becomes a real text box you can edit in
PowerPoint, and every picture (GIFs included) is placed as its own object,
so you can move it, resize it or swap the file.

    python3 tools/md2pptx.py

Output: state-of-openhistoricalmap-2026.pptx in the repo root.
"""

import io
import os
import re
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(ROOT, 'slides.md')
OUT = os.path.join(ROOT, 'state-of-openhistoricalmap-2026.pptx')

# --- paleta y tipografías del tema -------------------------------------------
INK = RGBColor(0x21, 0x1B, 0x13)
INK_SOFT = RGBColor(0x57, 0x4B, 0x39)
INK_FAINT = RGBColor(0x8A, 0x7A, 0x60)
ROSA = RGBColor(0xBC, 0x60, 0x76)
ROSA_DEEP = RGBColor(0x8C, 0x38, 0x4E)
ROSA_TINT = RGBColor(0xE7, 0xCB, 0xCD)
TEAL = RGBColor(0x2C, 0x8A, 0x7B)
PAPER_2 = RGBColor(0xFA, 0xF8, 0xF4)
PAPER_EDGE = RGBColor(0xE4, 0xDE, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_DISPLAY = 'Fraunces'      # cae en Georgia si no está instalada
F_BODY = 'Hanken Grotesk'   # cae en Helvetica / Arial
F_MONO = 'IBM Plex Mono'    # cae en Menlo / Consolas

W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = W - 2 * MARGIN


# --- markdown mínimo ----------------------------------------------------------
def inline_runs(text):
    """Parte un texto en trozos (texto, bold, code)."""
    text = re.sub(r'<br\s*/?>', '\n', text)
    text = re.sub(r'<a [^>]*>(.*?)</a>', r'\1', text)
    text = re.sub(r'<[^>]+>', '', text)
    out = []
    for part in re.split(r'(\*\*.+?\*\*|`[^`]+`|\*[^*]+\*)', text):
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            out.append((part[2:-2], 'bold'))
        elif part.startswith('`') and part.endswith('`'):
            out.append((part[1:-1], 'code'))
        elif part.startswith('*') and part.endswith('*') and len(part) > 2:
            out.append((part[1:-1], 'italic'))
        else:
            out.append((part, 'plain'))
    return out


def images_in(line):
    """Saca (ruta, alt) de ![alt](src) y de <img src=...>."""
    found = []
    for m in re.finditer(r'!\[([^\]]*)\]\(([^)]+)\)', line):
        found.append((m.group(2), m.group(1)))
    for m in re.finditer(r'<img[^>]*src="([^"]+)"[^>]*?(?:alt="([^"]*)")?[^>]*>', line):
        found.append((m.group(1), m.group(2) or ''))
    return found


def parse_slides(md):
    md = re.sub(r'^<!--.*?-->\s*', '', md, flags=re.S)
    slides = []
    for chunk in re.split(r'\n---\n', md):
        head = re.search(r'<!--\s*\.slide:\s*class="([^"]*)"(.*?)-->', chunk, flags=re.S)
        cls = head.group(1).split() if head else []
        attrs = dict(re.findall(r'data-([a-z-]+)="([^"]*)"', head.group(2))) if head else {}
        body = chunk[head.end():] if head else chunk
        note = ''
        m = re.search(r'^Note:\s*(.*)$', body, flags=re.M | re.S)
        if m:
            note = m.group(1).strip()
            body = body[:m.start()]
        lines = [l.rstrip() for l in body.strip('\n').split('\n')]
        slides.append({'class': cls, 'attrs': attrs, 'lines': lines, 'note': note})
    return slides


# --- helpers de dibujo --------------------------------------------------------
def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return box, tf


def write(tf, text, size, font=F_BODY, color=INK, bold=False, italic=False,
          space_after=0, align=None, line_spacing=None, first=False, caps=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        para.alignment = align
    para.space_after = Pt(space_after)
    if line_spacing:
        para.line_spacing = line_spacing
    for chunk, kind in inline_runs(text):
        for i, piece in enumerate(chunk.split('\n')):
            if i:
                para.add_line_break()
            run = para.add_run()
            run.text = piece.upper() if caps else piece
            f = run.font
            f.size = Pt(size)
            f.name = F_MONO if kind == 'code' else font
            f.bold = bold or kind == 'bold'
            f.italic = italic or kind == 'italic'
            f.color.rgb = ROSA_DEEP if kind == 'code' else color
    return para


def kicker(slide, text, y=Inches(0.5)):
    box, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3))
    write(tf, text, 12, F_MONO, ROSA, bold=True, first=True, caps=True)
    return box


def picture(slide, src, box_x, box_y, box_w, box_h):
    """Mete la imagen dentro de la caja, centrada, sin deformarla."""
    path = os.path.join(ROOT, src)
    if not os.path.exists(path):
        print('  missing image:', src)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = int(box_x + (box_w - w) / 2)
    y = int(box_y + (box_h - h) / 2)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAPER_2
    shape.line.color.rgb = PAPER_EDGE
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def pill(slide, text, x, y, size=14):
    w = Inches(0.22 * len(text) + 0.4)
    h = Inches(0.42)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.2
    shape.fill.solid()
    shape.fill.fore_color.rgb = ROSA_TINT
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, text, size, F_MONO, ROSA_DEEP, bold=True, first=True,
          align=PP_ALIGN.CENTER)
    return shape, w


# --- una función por tipo de slide -------------------------------------------
def build_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    cls = s['class']
    lines = [l for l in s['lines'] if l.strip()]
    heads = {'h1': [], 'h2': [], 'h3': []}
    paras, bullets, rows, quotes, imgs, pills = [], [], [], [], [], []

    for l in lines:
        st = l.strip()
        if st.startswith('### '):
            heads['h3'].append(st[4:])
        elif st.startswith('## '):
            heads['h2'].append(st[3:])
        elif st.startswith('# '):
            heads['h1'].append(st[2:])
        elif st.startswith('|'):
            cells = [c.strip() for c in st.strip('|').split('|')]
            if not all(re.fullmatch(r':?-{2,}:?', c) for c in cells):
                rows.append(cells)
        elif st.startswith('>'):
            quotes.append(st.lstrip('> ').strip())
        elif st.startswith('- '):
            bullets.append([st[2:].strip()])
        elif bullets and l.startswith('  '):
            bullets[-1].append(st)
        elif images_in(st):
            imgs += images_in(st)
        elif re.fullmatch(r'`[^`]+`', st):
            pills.append(st.strip('`'))
        elif st.startswith('<!--'):
            continue
        else:
            paras.append(st)

    y = Inches(0.5)
    if heads['h3']:
        kicker(slide, heads['h3'][0], y)
        y = Inches(0.95)

    def title(text, size, gap=0.35):
        nonlocal y
        nlines = text.count('<br>') + text.count('<br/>') + 1
        h = Inches(size / 72 * 1.2 * nlines)
        box, tf = textbox(slide, MARGIN, y, CONTENT_W, h)
        write(tf, text, size, F_DISPLAY, INK, first=True, line_spacing=1.1)
        y = y + h + Inches(gap)

    # --- portada
    if 'cover' in cls:
        if heads['h1']:
            title(heads['h1'][0], 46)
        y = y + Inches(0.25)
        for i, p in enumerate(paras):
            box, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.5))
            size = 20 if i == 0 else 13
            write(tf, p, size, F_BODY if i == 0 else F_MONO,
                  INK_SOFT if i == 0 else INK_FAINT, first=True)
            y = y + Inches(0.55 if i == 0 else 0.4)
        return slide

    # --- ponentes
    if 'speakers' in cls and rows:
        pics = images_in(' '.join(rows[0]))
        names = rows[1] if len(rows) > 1 else []
        roles = rows[2] if len(rows) > 2 else []
        n = max(len(pics), 1)
        colw = CONTENT_W / n
        for i in range(n):
            cx = MARGIN + colw * i
            if i < len(pics):
                picture(slide, pics[i][0], int(cx + colw / 2 - Inches(1.2)),
                        Inches(1.5), Inches(2.4), Inches(2.4))
            if i < len(names):
                box, tf = textbox(slide, cx, Inches(4.15), colw, Inches(0.5),
                                  align=PP_ALIGN.CENTER)
                write(tf, names[i], 26, F_DISPLAY, INK, first=True,
                      align=PP_ALIGN.CENTER)
            if i < len(roles):
                box, tf = textbox(slide, cx, Inches(4.8), colw, Inches(0.8),
                                  align=PP_ALIGN.CENTER)
                write(tf, roles[i], 12, F_MONO, INK_FAINT, first=True,
                      align=PP_ALIGN.CENTER)
        return slide

    # --- una animación o imagen grande, con el título encima
    if 'hero' in cls:
        if imgs:
            picture(slide, imgs[0][0], MARGIN, Inches(1.15), CONTENT_W, Inches(5.6))
        return slide

    # --- la ecuación
    if 'equation' in cls:
        head = heads['h2'][0] if heads['h2'] else ''
        parts = re.split(r'(`[^`]+`)', head)
        words = ''.join(p for p in parts if not p.startswith('`'))
        codes = [p.strip('`') for p in parts if p.startswith('`')]
        box, tf = textbox(slide, MARGIN, Inches(0.8), CONTENT_W, Inches(0.9),
                          align=PP_ALIGN.CENTER)
        write(tf, words.strip(), 34, F_DISPLAY, INK, first=True, align=PP_ALIGN.CENTER)
        widths = [Inches(0.22 * len(c) + 0.4) for c in codes]
        total = sum(widths) + Inches(0.2) * (len(codes) - 1)
        px = int(MARGIN + (CONTENT_W - total) / 2)
        for c, w in zip(codes, widths):
            pill(slide, c, px, Inches(1.85))
            px = px + w + Inches(0.2)
        if imgs:
            src = imgs[0][0].replace('.svg', '.png')
            picture(slide, src, MARGIN, Inches(2.6), CONTENT_W, Inches(4.2))
        return slide

    # --- tarjetas de datos
    if 'facts' in cls:
        if heads['h2']:
            title(heads['h2'][0], 34)
        cols = 2 if 'g2' in cls else 3
        n = len(bullets)
        rows_n = (n + cols - 1) // cols
        gap = Inches(0.22)
        cw = int((CONTENT_W - gap * (cols - 1)) / cols)
        ch = Inches(1.55) if rows_n <= 2 else Inches(1.2)
        top = y + Inches(0.2)
        for i, item in enumerate(bullets):
            cx = MARGIN + (cw + gap) * (i % cols)
            cy = top + (ch + gap) * (i // cols)
            card(slide, cx, cy, cw, ch)
            box, tf = textbox(slide, cx + Inches(0.28), cy + Inches(0.22),
                              cw - Inches(0.5), ch - Inches(0.4))
            write(tf, item[0], 34, F_DISPLAY, INK, first=True, space_after=4)
            if len(item) > 1:
                write(tf, item[1], 11, F_MONO, INK_FAINT, caps=True, space_after=3)
            if len(item) > 2:
                write(tf, ' '.join(item[2:]), 12, F_BODY, INK_SOFT)
        if paras:
            box, tf = textbox(slide, MARGIN, top + (ch + gap) * rows_n + Inches(0.1),
                              CONTENT_W, Inches(0.5))
            for i, p in enumerate(paras):
                write(tf, p, 11, F_MONO, INK_FAINT, first=(i == 0))
        return slide

    # --- imagen a sangre completa con pie
    if 'bleed' in cls:
        bg = s['attrs'].get('background-image')
        if bg:
            picture(slide, bg, Inches(0), Inches(0.2), W, Inches(6.4))
        if quotes:
            box, tf = textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.6),
                              align=PP_ALIGN.CENTER)
            for i, q in enumerate(quotes):
                write(tf, q, 13, F_MONO, INK_SOFT, first=(i == 0),
                      align=PP_ALIGN.CENTER)
        return slide

    # --- grilla de imágenes con pies
    if 'imgrid' in cls:
        if heads['h2']:
            title(heads['h2'][0], 34)
        if paras:
            box, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.4))
            write(tf, paras[0], 15, F_BODY, INK_SOFT, first=True)
            y = y + Inches(0.55)
        cells = images_in(' '.join(rows[0])) if rows else []
        caps = rows[1] if len(rows) > 1 else []
        n = max(len(cells), 1)
        gap = Inches(0.3)
        cw = int((CONTENT_W - gap * (n - 1)) / n)
        for i, (src, alt) in enumerate(cells):
            cx = MARGIN + (cw + gap) * i
            picture(slide, src, cx, y, cw, Inches(3.1))
            if i < len(caps):
                box, tf = textbox(slide, cx, y + Inches(3.25), cw, Inches(0.8))
                write(tf, caps[i], 11, F_MONO, INK_SOFT, first=True, line_spacing=1.3)
        if len(paras) > 1:
            box, tf = textbox(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.5))
            write(tf, paras[1], 13, F_BODY, INK_SOFT, first=True)
        return slide

    # --- captura a pantalla completa con su URL
    if 'shot' in cls:
        top = Inches(0.5) if not heads['h3'] else Inches(1.0)
        if imgs:
            picture(slide, imgs[0][0], MARGIN, top, CONTENT_W, Inches(5.5))
        if pills:
            box, tf = textbox(slide, MARGIN, Inches(6.35), CONTENT_W, Inches(0.4),
                              align=PP_ALIGN.CENTER)
            write(tf, pills[0], 14, F_MONO, INK_SOFT, first=True, align=PP_ALIGN.CENTER)
        return slide

    # --- cita grande
    if 'quote' in cls:
        box, tf = textbox(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(2.5))
        write(tf, quotes[0], 34, F_DISPLAY, INK, italic=True, first=True,
              line_spacing=1.15, space_after=18)
        for q in quotes[1:]:
            write(tf, q, 12, F_MONO, INK_FAINT, caps=True, line_spacing=1.4)
        return slide

    # --- dos columnas
    if 'cols' in cls:
        if heads['h1']:
            title(heads['h1'][0], 40)
        y = y + Inches(0.3)
        colw = int((CONTENT_W - Inches(0.6)) / 2)
        groups = []
        cur = None
        for l in lines:
            st = l.strip()
            if st.startswith('### '):
                cur = {'head': st[4:], 'items': []}
                groups.append(cur)
            elif st.startswith('- ') and cur:
                cur['items'].append(st[2:])
        for i, g in enumerate(groups[:2]):
            cx = MARGIN + (colw + Inches(0.6)) * i
            box, tf = textbox(slide, cx, y, colw, Inches(0.5))
            write(tf, g['head'], 22, F_DISPLAY, ROSA_DEEP, first=True)
            box, tf = textbox(slide, cx, y + Inches(0.65), colw, Inches(2.6))
            for j, it in enumerate(g['items']):
                write(tf, it, 15, F_BODY, INK_SOFT, first=(j == 0), space_after=8,
                      line_spacing=1.25)
        if paras:
            box, tf = textbox(slide, MARGIN, Inches(6.4), CONTENT_W, Inches(0.5))
            write(tf, paras[-1], 11, F_MONO, INK_FAINT, first=True)
        return slide

    # --- cierre
    if 'closing' in cls:
        if heads['h1']:
            box, tf = textbox(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(1.0),
                              align=PP_ALIGN.CENTER)
            write(tf, heads['h1'][0], 42, F_DISPLAY, INK, first=True,
                  align=PP_ALIGN.CENTER)
        labels = re.findall(r'<span class="qr-label">(.*?)</span>', ' '.join(lines))
        n = max(len(labels), 1)
        colw = CONTENT_W / n
        for i, lab in enumerate(labels):
            box, tf = textbox(slide, MARGIN + colw * i, Inches(3.4), colw, Inches(1.0),
                              align=PP_ALIGN.CENTER)
            write(tf, re.sub(r'<br\s*/?>', '\n', lab), 12, F_MONO, INK_SOFT,
                  first=True, align=PP_ALIGN.CENTER, line_spacing=1.4)
        tail = [p for p in paras if not p.startswith('<div')]
        if tail:
            box, tf = textbox(slide, MARGIN, Inches(5.6), CONTENT_W, Inches(0.6),
                              align=PP_ALIGN.CENTER)
            write(tf, tail[-1], 13, F_MONO, INK_SOFT, first=True, align=PP_ALIGN.CENTER)
        return slide

    # --- cualquier otra: título, párrafos, pastillas
    if heads['h1']:
        title(heads['h1'][0], 40)
    if heads['h2']:
        title(heads['h2'][0], 32)
    for p in paras:
        box, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.6))
        write(tf, p, 17, F_BODY, INK_SOFT, first=True, line_spacing=1.3)
        y = y + Inches(0.7)
    px = MARGIN
    for c in pills:
        shape, w = pill(slide, c, px, y)
        px = px + w + Inches(0.15)
    return slide


def main():
    md = io.open(SRC, encoding='utf-8').read()
    slides = parse_slides(md)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for i, s in enumerate(slides):
        print(f'slide {i + 1:2d}  {" ".join(s["class"])}')
        slide = build_slide(prs, s)
        if s['note']:
            slide.notes_slide.notes_text_frame.text = s['note']
    prs.save(OUT)
    print('saved', OUT, f'({len(slides)} slides)')


if __name__ == '__main__':
    main()
