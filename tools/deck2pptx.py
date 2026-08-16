#!/usr/bin/env python3
"""
Builds a PowerPoint deck that mirrors the web slides, keeping the text editable.

It does not screenshot anything: it reads a layout.json produced by
tools/measure_deck.mjs (which measures every text box, font size, colour and
picture straight from the browser) and rebuilds the same boxes in PowerPoint.
GIFs go in as GIFs, so they still animate.

    python3 -m http.server 8850 &
    node tools/measure_deck.mjs        # writes tools/layout.json
    python3 tools/deck2pptx.py

Output: state-of-openhistoricalmap-2026.pptx in the repo root.
"""

import io
import json
import os
import re

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LAYOUT = os.path.join(ROOT, 'tools', 'layout.json')
SRC = os.path.join(ROOT, 'slides.md')
OUT = os.path.join(ROOT, 'state-of-openhistoricalmap-2026.pptx')

STAGE_W, STAGE_H = 1280.0, 720.0
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
PX = SLIDE_W / STAGE_W          # EMU por píxel de la web
PT_PER_PX = 0.75                # 96 dpi → puntos

ALIGN = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
         'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY,
         'start': PP_ALIGN.LEFT, 'end': PP_ALIGN.RIGHT}


def px(v):
    return int(round(v * PX))


def color(css):
    m = re.match(r'rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)', css or '')
    if not m:
        return RGBColor(0x21, 0x1B, 0x13)
    return RGBColor(*(int(float(x)) for x in m.groups()))


def family(css):
    first = (css or '').split(',')[0].strip().strip('"\'')
    return first or 'Helvetica'


def notes_from_md():
    md = io.open(SRC, encoding='utf-8').read()
    md = re.sub(r'^<!--.*?-->\s*', '', md, flags=re.S)
    out = []
    for chunk in re.split(r'\n---\n', md):
        m = re.search(r'^Note:\s*(.*)$', chunk, flags=re.M | re.S)
        out.append(m.group(1).strip() if m else '')
    return out


def add_text(slide, item):
    box = slide.shapes.add_textbox(px(item['x']), px(item['y']),
                                   px(item['w']) + Inches(0.06), px(item['h']) + Inches(0.12))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    para = tf.paragraphs[0]
    para.alignment = ALIGN.get(item.get('align', 'left'), PP_ALIGN.LEFT)
    # interlineado como múltiplo: el exacto en puntos desplaza el texto hacia arriba
    base = max((r['size'] for r in item['runs']), default=16)
    if item.get('lineHeight') and base:
        para.line_spacing = max(item['lineHeight'] / base, 0.9)

    for run in item['runs']:
        text = run['text']
        if run.get('transform') == 'uppercase':
            text = text.upper()
        for i, piece in enumerate(text.split('\n')):
            if i:
                para.add_line_break()
            if not piece:
                continue
            r = para.add_run()
            r.text = piece
            f = r.font
            f.size = Pt(max(run['size'] * PT_PER_PX, 6))
            f.name = family(run['family'])
            f.bold = run.get('weight', 400) >= 600
            f.italic = bool(run.get('italic'))
            f.color.rgb = color(run.get('color'))
    return box


def add_box(slide, item):
    """Pastillas, tarjetas y pies de foto: el rectángulo que va detrás del texto."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if item.get('radius', 0) > 1 else MSO_SHAPE.RECTANGLE,
        px(item['x']), px(item['y']), px(item['w']), px(item['h']))
    if item.get('radius', 0) > 1:
        shape.adjustments[0] = min(0.5, item['radius'] / max(item['h'], 1))
    if item.get('alpha', 0) > 0.02:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(item['color'])
    else:
        shape.fill.background()
    if item.get('border'):
        shape.line.color.rgb = color(item['border'])
        shape.line.width = Pt(max(item.get('borderWidth', 1) * PT_PER_PX, 0.5))
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_picture(slide, src, x, y, w, h):
    rel = src.split('8850/')[-1] if '8850/' in src else src
    rel = rel.split('?')[0].lstrip('/')
    if rel.endswith('.svg') and os.path.exists(os.path.join(ROOT, rel[:-4] + '.png')):
        rel = rel[:-4] + '.png'          # PowerPoint no dibuja SVG del tema
    path = os.path.join(ROOT, rel)
    if not os.path.exists(path):
        print('   missing:', rel)
        return None
    return slide.shapes.add_picture(path, px(x), px(y), px(w), px(h))


def add_background_picture(slide, src):
    """Los bleed llevan el GIF como fondo con background-size: contain."""
    rel = src.split('8850/')[-1].split('?')[0].lstrip('/')
    path = os.path.join(ROOT, rel)
    if not os.path.exists(path):
        print('   missing background:', rel)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(STAGE_W / iw, STAGE_H / ih)
    w, h = iw * scale, ih * scale
    add_picture(slide, rel, (STAGE_W - w) / 2, (STAGE_H - h) / 2, w, h)


def main():
    layout = json.load(io.open(LAYOUT, encoding='utf-8'))
    notes = notes_from_md()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    for s in layout:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color(s.get('bgColor'))

        if s.get('bgImage'):
            add_background_picture(slide, s['bgImage'])

        for item in s['items']:
            if item['kind'] == 'box':
                add_box(slide, item)
            elif item['kind'] == 'img':
                add_picture(slide, item['src'], item['x'], item['y'], item['w'], item['h'])
            else:
                add_text(slide, item)

        i = s['index']
        if i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]
        print(f'slide {i + 1:2d}  {len(s["items"])} items'
              f'{"  + background" if s.get("bgImage") else ""}')

    prs.save(OUT)
    print('saved', OUT, f'({len(layout)} slides)')


if __name__ == '__main__':
    main()
